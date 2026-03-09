#!/usr/bin/env python3
"""
Script para gerar apresentação de slides PowerPoint da Apostila de Git
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_title_slide(prs, title, subtitle):
    """Cria slide de título"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    # Título principal
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Subtítulo
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle
    
    return slide

def create_content_slide(prs, title, content_items):
    """Cria slide de conteúdo com bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # Título
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Conteúdo
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.level = 0
        
        # Se for código, mudar formatação
        if item.startswith('  '):
            p.font.name = 'Courier New'
            p.font.size = Pt(16)
    
    return slide

def create_two_column_slide(prs, title, left_content, right_content):
    """Cria slide com duas colunas"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Coluna esquerda
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.3), Inches(5))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    for i, item in enumerate(left_content):
        if i == 0:
            p = tf_left.paragraphs[0]
        else:
            p = tf_left.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
    
    # Coluna direita
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(5))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    for i, item in enumerate(right_content):
        if i == 0:
            p = tf_right.paragraphs[0]
        else:
            p = tf_right.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
    
    return slide

def create_code_slide(prs, title, code, explanation=None):
    """Cria slide com código"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    
    # Caixa de código
    code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(3.5))
    tf_code = code_box.text_frame
    tf_code.word_wrap = True
    
    p = tf_code.paragraphs[0]
    p.text = code
    p.font.name = 'Courier New'
    p.font.size = Pt(14)
    
    # Explicação
    if explanation:
        exp_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.7), Inches(9), Inches(1.5))
        tf_exp = exp_box.text_frame
        tf_exp.word_wrap = True
        p = tf_exp.paragraphs[0]
        p.text = explanation
        p.font.size = Pt(16)
    
    return slide

def create_section_slide(prs, section_number, section_title):
    """Cria slide de seção/divisão"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Fundo colorido (simulado com forma)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(45, 55, 72)
    shape.line.fill.background()
    
    # Número da seção
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(2), Inches(1.5))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_number
    p.font.size = Pt(72)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Título da seção
    title_box = slide.shapes.add_textbox(Inches(2.5), Inches(2.5), Inches(7), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(40)
    p.font.bold = True
    
    return slide

def main():
    prs = Presentation()
    
    # Configurar tamanho do slide (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # =========================================================================
    # SLIDE 1: Capa
    # =========================================================================
    slide = create_title_slide(
        prs,
        "Git e GitHub",
        "Versionamento de Software do Básico ao Avançado\n\nApostila Completa"
    )
    
    # =========================================================================
    # SLIDE 2: Agenda
    # =========================================================================
    agenda = [
        "1. Introdução ao Versionamento",
        "2. Git - Fundamentos e Comandos Básicos",
        "3. GitHub - Colaboração e Hospedagem",
        "4. Tópicos Avançados",
        "5. Exercícios Práticos"
    ]
    create_content_slide(prs, "Agenda", agenda)
    
    # =========================================================================
    # SEÇÃO 1: INTRODUÇÃO
    # =========================================================================
    create_section_slide(prs, "01", "Introdução ao Versionamento")
    
    # =========================================================================
    # SLIDE 3: O que é Versionamento?
    # =========================================================================
    versionamento = [
        "Processo de rastrear e gerenciar alterações em um projeto",
        "",
        "Permite:",
        "• Documentar todas as mudanças",
        "• Reverter para versões anteriores",
        "• Trabalhar em equipe sem conflitos",
        "• Manter backup completo do projeto"
    ]
    create_content_slide(prs, "O que é Versionamento de Software?", versionamento)
    
    # =========================================================================
    # SLIDE 4: Por que usar Versionamento?
    # =========================================================================
    slide = create_two_column_slide(
        prs,
        "Por que usar Versionamento?",
        # Left - Problemas sem versionamento
        [
            "SEM VERSIONAMENTO:",
            "",
            "❌ arquivo_final.zip",
            "❌ arquivo_final_v2.zip",
            "❌ arquivo_final_agora_vai.zip",
            "❌ arquivo_final_seria_esse.zip",
            "",
            "• Perda de dados",
            "• Sem histórico de mudanças",
            "• Conflitos na equipe"
        ],
        # Right - Benefícios
        [
            "COM VERSIONAMENTO:",
            "",
            "✅ Histórico completo",
            "✅ Cada mudança documentada",
            "✅ Quem, quando e porquê",
            "✅ Reversão segura",
            "✅ Trabalho em paralelo",
            "✅ Backup automático"
        ]
    )
    
    # =========================================================================
    # SLIDE 5: Sistemas de Versionamento
    # =========================================================================
    slide = create_two_column_slide(
        prs,
        "Tipos de Sistemas",
        # Centralizado
        [
            "CENTRALIZADO (CVCS)",
            "",
            "• Servidor único central",
            "• Clientes se conectam",
            "• Exemplos: SVN, CVS",
            "",
            "Desvantagens:",
            "• Ponto único de falha",
            "• Requer conexão constante"
        ],
        # Distribuído
        [
            "DISTRIBUÍDO (DVCS)",
            "",
            "• Cada clone é backup completo",
            "• Trabalhe offline",
            "• Exemplos: Git, Mercurial",
            "",
            "Vantagens:",
            "• Resiliência",
            "• Performance",
            "• Flexibilidade"
        ]
    )
    
    # =========================================================================
    # SEÇÃO 2: GIT
    # =========================================================================
    create_section_slide(prs, "02", "Git - Fundamentos")
    
    # =========================================================================
    # SLIDE 6: O que é Git?
    # =========================================================================
    git_info = [
        "Sistema de Controle de Versão Distribuído",
        "",
        "Criado por Linus Torvalds (2005)",
        "",
        "Características:",
        "• Gratuito e Open Source",
        "• Rápido e eficiente",
        "• Seguro (SHA-1 hashing)",
        "• Usado mundialmente"
    ]
    create_content_slide(prs, "O que é Git?", git_info)
    
    # =========================================================================
    # SLIDE 7: Ciclo de Vida dos Arquivos
    # =========================================================================
    slide = create_two_column_slide(
        prs,
        "Ciclo de Vida dos Arquivos",
        # Estados
        [
            "ESTADOS:",
            "",
            "1. Untracked (Não rastreado)",
            "   • Arquivo novo",
            "",
            "2. Modified (Modificado)",
            "   • Arquivo alterado",
            "",
            "3. Staged (Preparado)",
            "   • Pronto para commit",
            "",
            "4. Committed (Salvo)",
            "   • No histórico Git"
        ],
        # Comandos
        [
            "COMANDOS:",
            "",
            "git add <arquivo>",
            "  → Untracked → Staged",
            "",
            "git commit",
            "  → Staged → Committed",
            "",
            "git status",
            "  → Verifica estado",
            "",
            "git restore",
            "  → Descarta mudanças"
        ]
    )
    
    # =========================================================================
    # SLIDE 8: Instalação
    # =========================================================================
    instalacao = [
        "Windows:",
        "  • Baixe em git-scm.com",
        "  • Execute o instalador",
        "",
        "macOS:",
        "  • brew install git",
        "  • ou xcode-select --install",
        "",
        "Linux:",
        "  • sudo apt install git (Debian/Ubuntu)",
        "  • sudo dnf install git (Fedora)",
        "",
        "Verificar:",
        "  git --version"
    ]
    create_content_slide(prs, "Instalação do Git", instalacao)
    
    # =========================================================================
    # SLIDE 9: Configuração Inicial
    # =========================================================================
    config = """git config --global user.name "Seu Nome"
git config --global user.email "seu@email.com"
git config --global core.editor "code --wait"

# Verificar configurações
git config --list"""
    
    create_code_slide(
        prs,
        "Configuração Inicial",
        config,
        "Configure nome, e-mail e editor antes de começar!"
    )
    
    # =========================================================================
    # SLIDE 10: Comandos Básicos - Parte 1
    # =========================================================================
    comandos1 = [
        "INICIAR REPOSITÓRIO:",
        "  git init                    # Novo repositório",
        "  git clone <url>             # Copiar repositório",
        "",
        "STATUS E LOG:",
        "  git status                  # Ver estado dos arquivos",
        "  git log                     # Histórico de commits",
        "  git log --oneline           # Histórico resumido",
        "",
        "ADICIONAR ARQUIVOS:",
        "  git add <arquivo>           # Adiciona arquivo específico",
        "  git add .                   # Adiciona todos"
    ]
    create_content_slide(prs, "Comandos Básicos - Parte 1", comandos1)
    
    # =========================================================================
    # SLIDE 11: Comandos Básicos - Parte 2
    # =========================================================================
    comandos2 = [
        "COMMITS:",
        "  git commit -m \"mensagem\"    # Cria commit",
        "  git commit -am \"msg\"        # Add + commit",
        "",
        "REMOTO:",
        "  git remote add origin <url>   # Adiciona remoto",
        "  git push origin main          # Envia alterações",
        "  git pull origin main          # Baixa alterações",
        "",
        "BRANCHES:",
        "  git branch                    # Lista branches",
        "  git checkout -b <nome>        # Cria e muda",
        "  git checkout <nome>           # Muda de branch"
    ]
    create_content_slide(prs, "Comandos Básicos - Parte 2", comandos2)
    
    # =========================================================================
    # SLIDE 12: Exemplo Prático
    # =========================================================================
    exemplo = """# Criar projeto
mkdir meu-projeto && cd meu-projeto
git init

# Criar arquivo
echo "# Meu Projeto" > README.md
git add README.md
git commit -m "feat: adiciona README"

# Conectar com GitHub
git remote add origin https://github.com/user/repo.git
git push -u origin main"""
    
    create_code_slide(
        prs,
        "Exemplo: Primeiro Repositório",
        exemplo,
        "Fluxo completo do zero até o GitHub"
    )
    
    # =========================================================================
    # SEÇÃO 3: GITHUB
    # =========================================================================
    create_section_slide(prs, "03", "GitHub")
    
    # =========================================================================
    # SLIDE 13: O que é GitHub?
    # =========================================================================
    github = [
        "Plataforma de hospedagem de repositórios Git",
        "",
        "Funcionalidades:",
        "• Repositórios públicos e privados",
        "• Pull Requests para code review",
        "• Issues para gerenciamento de tarefas",
        "• GitHub Actions (CI/CD)",
        "• GitHub Pages (sites estáticos)",
        "• 100+ milhões de desenvolvedores"
    ]
    create_content_slide(prs, "O que é GitHub?", github)
    
    # =========================================================================
    # SLIDE 14: Autenticação
    # =========================================================================
    auth = [
        "SSH (Recomendado):",
        "  ssh-keygen -t ed25519",
        "  # Adicione ~/.ssh/id_ed25519.pub no GitHub",
        "  git clone git@github.com:user/repo.git",
        "",
        "Personal Access Token:",
        "  • Settings → Developer settings",
        "  • Generate new token (classic)",
        "  • Selecione permissão 'repo'",
        "  • Use token no lugar da senha",
        "",
        "⚠️ Senha não funciona mais para Git!"
    ]
    create_content_slide(prs, "Autenticação no GitHub", auth)
    
    # =========================================================================
    # SLIDE 15: Pull Requests
    # =========================================================================
    pr_flow = [
        "FLUXO DO PULL REQUEST:",
        "",
        "1. Fork do repositório (se necessário)",
        "2. git checkout -b feature/nova",
        "3. Faça alterações e commit",
        "4. git push origin feature/nova",
        "5. Abra Pull Request no GitHub",
        "6. Aguarde code review",
        "7. Merge após aprovação",
        "",
        "PR permite colaboração segura!"
    ]
    create_content_slide(prs, "Pull Requests", pr_flow)
    
    # =========================================================================
    # SLIDE 16: Issues
    # =========================================================================
    issues = [
        "Gerencie bugs e tarefas",
        "",
        "Labels comuns:",
        "  🐛 bug - Algo quebrado",
        "  ✨ enhancement - Nova feature",
        "  📚 documentation - Docs",
        "  🟢 good first issue - Iniciantes",
        "",
        "Boas práticas:",
        "  • Título descritivo",
        "  • Passos para reproduzir",
        "  • Comportamento esperado",
        "  • Screenshots quando útil"
    ]
    create_content_slide(prs, "Issues - Gerenciamento de Tarefas", issues)
    
    # =========================================================================
    # SEÇÃO 4: AVANÇADO
    # =========================================================================
    create_section_slide(prs, "04", "Tópicos Avançados")
    
    # =========================================================================
    # SLIDE 17: Branches
    # =========================================================================
    branches = [
        "Convenção de nomes:",
        "  feature/login           # Nova funcionalidade",
        "  bugfix/erro-calculo     # Correção de bug",
        "  hotfix/seguranca        # Correção urgente",
        "  docs/readme             # Documentação",
        "  refactor/modulo-x       # Refatoração",
        "",
        "Comandos:",
        "  git branch -a           # Lista todas",
        "  git checkout -b <nome>  # Cria e muda",
        "  git merge <branch>      # Mescla",
        "  git branch -d <nome>    # Deleta"
    ]
    create_content_slide(prs, "Branches - Ramificações", branches)
    
    # =========================================================================
    # SLIDE 18: Merge vs Rebase
    # =========================================================================
    slide = create_two_column_slide(
        prs,
        "Merge vs Rebase",
        # Merge
        [
            "MERGE",
            "",
            "• Preserva histórico",
            "• Cria commit de merge",
            "• Mais seguro",
            "",
            "git merge feature",
            "",
            "Use em:",
            "• Branches públicas",
            "• Pull requests"
        ],
        # Rebase
        [
            "REBASE",
            "",
            "• Reescreve histórico",
            "• Histórico linear",
            "• Mais limpo",
            "",
            "git rebase main",
            "",
            "Use em:",
            "• Branches locais",
            "• Antes de PR",
            "",
            "⚠️ Nunca em branches públicas!"
        ]
    )
    
    # =========================================================================
    # SLIDE 19: Stash
    # =========================================================================
    stash = [
        "Guarda alterações temporariamente",
        "",
        "Comandos:",
        "  git stash                    # Guarda mudanças",
        "  git stash save \"msg\"         # Com mensagem",
        "  git stash list               # Lista stashes",
        "  git stash pop                # Restaura e remove",
        "  git stash apply              # Restaura (mantém)",
        "  git stash drop               # Remove stash",
        "",
        "Cenário: Mudar de branch sem commitar!"
    ]
    create_content_slide(prs, "Stash - Alterações Temporárias", stash)
    
    # =========================================================================
    # SLIDE 20: Tags
    # =========================================================================
    tags = [
        "Marca versões específicas",
        "",
        "Versionamento Semântico:",
        "  MAIOR.MENOR.CORREÇÃO",
        "  1.0.0  → Lançamento inicial",
        "  1.0.1  → Correção de bug",
        "  1.1.0  → Nova funcionalidade",
        "  2.0.0  → Mudança incompatível",
        "",
        "Comandos:",
        "  git tag -a v1.0.0 -m \"msg\"   # Cria tag",
        "  git push origin --tags         # Envia tags"
    ]
    create_content_slide(prs, "Tags - Versionamento", tags)
    
    # =========================================================================
    # SLIDE 21: Reset, Restore, Revert
    # =========================================================================
    slide = create_two_column_slide(
        prs,
        "Desfazendo Alterações",
        # Reset
        [
            "GIT RESET",
            "",
            "Move HEAD, desfaz commits",
            "",
            "git reset --soft HEAD~1",
            "  → Mantém alterações",
            "",
            "git reset --hard HEAD~1",
            "  → Descarta tudo",
            "",
            "⚠️ Apenas local!"
        ],
        # Revert
        [
            "GIT REVERT",
            "",
            "Cria commit que desfaz",
            "",
            "git revert HEAD",
            "  → Reverte último",
            "",
            "git revert abc123",
            "  → Reverte específico",
            "",
            "✅ Seguro para remoto!"
        ]
    )
    
    # =========================================================================
    # SLIDE 22: Reflog
    # =========================================================================
    reflog = [
        "Seu seguro contra acidentes!",
        "",
        "Registra todos os movimentos do HEAD",
        "",
        "Comandos:",
        "  git reflog                 # Ver histórico",
        "  git reset --hard HEAD@{5}  # Volta estado",
        "",
        "Recuperar branch deletada:",
        "  1. git reflog",
        "  2. Ache o commit",
        "  3. git checkout -b nova-branch <hash>",
        "",
        "Nada se perde com reflog!"
    ]
    create_content_slide(prs, "Reflog - Recuperação de Desastres", reflog)
    
    # =========================================================================
    # SLIDE 23: Boas Práticas de Commit
    # =========================================================================
    commits = [
        "Mensagens ruins:",
        "  ❌ \"arrumei\"",
        "  ❌ \"teste\"",
        "  ❌ \"fix\"",
        "",
        "Mensagens boas:",
        "  ✅ \"Corrige cálculo de imposto\"",
        "  ✅ \"Adiciona validação de CPF\"",
        "",
        "Conventional Commits:",
        "  feat: Nova funcionalidade",
        "  fix: Correção de bug",
        "  docs: Documentação",
        "  refactor: Refatoração",
        "  test: Testes"
    ]
    create_content_slide(prs, "Boas Práticas de Commit", commits)
    
    # =========================================================================
    # SLIDE 24: Fluxo de Trabalho Recomendado
    # =========================================================================
    fluxo = [
        "1. git checkout main",
        "   git pull origin main",
        "",
        "2. git checkout -b feature/nova",
        "",
        "3. Trabalhe e faça commits atômicos",
        "",
        "4. git rebase main (opcional)",
        "",
        "5. git push -u origin feature/nova",
        "",
        "6. Crie Pull Request",
        "",
        "7. Após review: merge"
    ]
    create_content_slide(prs, "Fluxo de Trabalho Recomendado", fluxo)
    
    # =========================================================================
    # SEÇÃO 5: EXERCÍCIOS
    # =========================================================================
    create_section_slide(prs, "05", "Exercícios Práticos")
    
    # =========================================================================
    # SLIDE 25: Exercícios
    # =========================================================================
    exercicios = [
        "1. Configurar Git (nome, e-mail)",
        "",
        "2. Criar primeiro repositório",
        "",
        "3. Trabalhar com branches",
        "",
        "4. Simular e resolver conflitos",
        "",
        "5. Usar stash",
        "",
        "6. Criar tags de versão",
        "",
        "7. Pull Request no GitHub",
        "",
        "8. Fork e contribuição open source",
        "",
        "9. Recuperar com reflog",
        "",
        "10. Desafio final integrador"
    ]
    create_content_slide(prs, "Exercícios Práticos", exercicios)
    
    # =========================================================================
    # SLIDE 26: Recursos Adicionais
    # =========================================================================
    recursos = [
        "Documentação:",
        "  • git-scm.com/doc",
        "  • docs.github.com",
        "",
        "Prática:",
        "  • learngitbranching.js.org",
        "  • ohmygit.org",
        "",
        "Ferramentas:",
        "  • GitKraken, Sourcetree",
        "  • GitLens (VS Code)",
        "  • Meld (merge tool)",
        "",
        "Referência:",
        "  • Conventional Commits"
    ]
    create_content_slide(prs, "Recursos Adicionais", recursos)
    
    # =========================================================================
    # SLIDE 27: Obrigado
    # =========================================================================
    slide = create_title_slide(
        prs,
        "Obrigado!",
        "Dúvidas?\n\nBons estudos e bom código! 🚀"
    )
    
    # Salvar apresentação
    prs.save('apostila_git_slides.pptx')
    print("✅ Apresentação criada com sucesso: apostila_git_slides.pptx")
    print(f"📊 Total de slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
